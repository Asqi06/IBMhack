from pptx import Presentation
prs = Presentation(r'F:\ibm hackathon\_pptx_x\compiled.pptx')
print('slides:', len(prs.slides))
for i, slide in enumerate(prs.slides):
    title = slide.shapes.title.text if slide.shapes.title else '(no title)'
    body_parts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            txt = shape.text_frame.text.strip()
            if txt:
                body_parts.append(txt[:100])
    body = ' | '.join(body_parts) if body_parts else '(no text)'
    print(f'Slide {i+1}: title="{title}" body="{body}"')
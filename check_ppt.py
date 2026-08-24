from pptx import Presentation
prs = Presentation(r'F:\ibm hackathon\_pptx_x\compiled.pptx')
print('slides:', len(prs.slides))
for i, slide in enumerate(prs.slides):
    title = slide.shapes.title.text if slide.shapes.title else '(no title)'
    body = ''
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            body = shape.text_frame.text.strip()[:80]
            break
    print(f'Slide {i+1}: title="{title}" body="{body}"')
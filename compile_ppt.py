import os
import xml.etree.ElementTree as ET
from pptx import Presentation

ns = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
}

def extract_text(txBody):
    if txBody is None:
        return ''
    texts = txBody.findall('.//a:t', ns)
    if texts:
        return ''.join((t.text or '') for t in texts)
    return txBody.text or ''

slide_dir = r'F:/ibm hackathon/_pptx_x/ppt/slides'
slide_files = sorted([f for f in os.listdir(slide_dir) if f.startswith('slide') and f.endswith('.xml')])

prs = Presentation()
# use the second layout (index 1) which is title and content
for sf in slide_files:
    tree = ET.parse(os.path.join(slide_dir, sf))
    root = tree.getroot()
    title_text = ''
    body_text = ''
    for sp in root.findall('.//p:sp', ns):
        # find placeholder type
        ph = sp.find('p:nvSpPr/p:nvPr/p:ph', ns)
        ph_type = ph.get('type') if ph is not None else ''
        txBody = sp.find('p:txBody', ns)
        if ph_type == 'title':
            title_text = extract_text(txBody)
        elif ph_type == 'body':
            body_text = extract_text(txBody)
    # add slide with title and content layout
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title & Content
    slide.shapes.title.text = title_text[:200]
    # find the content text placeholder and set its text
    for shape in slide.shapes:
        if shape.has_text_frame:
            # if it's the placeholder text "Click to add text", replace it
            if shape.text_frame.text.startswith('Click to add'):
                shape.text_frame.clear()
                shape.text_frame.text = body_text[:500]
                break
    print(f"Processed {sf}: title={title_text[:30]}... body={body_text[:30]}...")

output_path = r'F:/ibm hackathon/_pptx_x/compiled.pptx'
prs.save(output_path)
print(f'Saved compiled PPTX to {output_path}')